"""import os
import google.generativeai as genai
from flask import Flask, render_template, request, send_file, jsonify
from pptx import Presentation

# Configure API Key
genai.configure(api_key="AIzaSyAOKA2pHefOM4hEsG-qphJLrq3uqcKDB0I")
model_name = "models/gemini-1.5-pro-latest"

app = Flask(__name__, static_url_path='/static')

# TEAM MEMBERS (edit this if needed)
TEAM_MEMBERS = ["Gourav Jat", "John Doe", "Jane Smith"]

def generate_slide_content(topic, num_slides=5, detail_level="average", style="classic"):
    detail_text = {
        "average": "moderate depth with bullet points or short paragraphs.",
        "detailed": "thorough and technical with bullet points or short paragraphs depending on context."
    }

    style_text = {
        "classic": "The presentation should have a classic, formal tone.",
        "modern": "The presentation should have a modern and visually appealing tone.",
        "dark": "The presentation should have a high-contrast, dark-themed tone."
    }

    # Prompt to Gemini
    prompt = (
        f"You are an expert in creating PowerPoint presentations.\n\n"
        f"Create a presentation on the topic: '{topic}'. The presentation should have {num_slides} slides.\n"
        f"Each slide must include a heading and either 5–9 detailed bullet points OR a short, well-written paragraph.\n"
        f"Ensure NO slide is empty or repeated.\n"
        f"The first slide should be a title slide and list the topic and the team members: {', '.join(TEAM_MEMBERS)}.\n"
        f"The remaining slides should cover subtopics and include meaningful, unique content.\n"
        f"The presentation should be {detail_text[detail_level]} and styled as follows: {style_text[style]}.\n\n"
        f"Return in this exact format:\n\n"
        f"Slide Heading: [Title of Slide]\n- Bullet 1 or paragraph sentence.\n- Bullet 2 or next sentence...\n\n"
        f"Repeat this format for each slide."
    )

    model = genai.GenerativeModel(model_name)
    try:
        response = model.generate_content(prompt)
        if not response or not hasattr(response, "text") or not response.text.strip():
            return [{"title": "Error", "bullets": ["No content generated. Please try again."]}]

        slides_raw = response.text.strip().split("\n\n")
        slides = []

        for slide in slides_raw:
            lines = slide.strip().split("\n")
            if not lines:
                continue
            title_line = lines[0].replace("Slide Heading:", "").strip()
            bullet_points = [line.replace("-", "").strip() for line in lines[1:] if line.startswith("-")]

            if not title_line or not bullet_points:
                continue  # Skip weak slides

            slides.append({
                "title": title_line,
                "bullets": bullet_points
            })

        # Enforce number of slides and fallback
        while len(slides) < num_slides:
            slides.append({
                "title": f"Extra Slide {len(slides) + 1}",
                "bullets": [f"This slide expands on the topic '{topic}' in detail."]
            })

        return slides[:num_slides]

    except Exception as e:
        return [{"title": "Error", "bullets": [str(e)]}]

def create_ppt(slide_contents, filename="AI_Presentation.pptx"):
    save_path = os.path.join(os.getcwd(), filename)
    ppt = Presentation()

    for slide_content in slide_contents:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        slide.shapes.title.text = slide_content["title"]
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()

        for point in slide_content["bullets"]:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0

    ppt.save(save_path)
    return save_path

@app.route("/")
def home():
    return render_template("index.html")

@app.route("/generate", methods=["POST"])
def generate_ppt():
    topic = request.form.get("topic")
    num_slides_raw = request.form.get("num_slides")
    if not num_slides_raw or not num_slides_raw.isdigit():
        return jsonify({"error": "Please enter a valid number of slides."}), 400

    num_slides = int(num_slides_raw)
    detail_level = request.form.get("detail_level", "average")
    style = request.form.get("style", "classic")

    num_slides = max(1, min(250, num_slides))  # Clamp to 1–250

    slides_preview = generate_slide_content(topic, num_slides, detail_level, style)

    if slides_preview and slides_preview[0].get("title") == "Error":
        return jsonify({"error": slides_preview[0]["bullets"][0]}), 500

    ppt_file_path = create_ppt(slides_preview)

    return jsonify({
        "slides": slides_preview,
        "ppt_file": True
    })

@app.route("/download", methods=["GET"])
def download_ppt():
    file_path = os.path.join(os.getcwd(), "AI_Presentation.pptx")
    if not os.path.exists(file_path):
        return jsonify({"error": "PPT file not found!"}), 404
    return send_file(file_path, as_attachment=True)

if __name__ == "__main__":
    app.run(debug=True) """


import os
import google.generativeai as genai
from flask import Flask, render_template, request, send_file, jsonify
from pptx import Presentation

# Configure API Key
genai.configure(api_key="AIzaSyAOKA2pHefOM4hEsG-qphJLrq3uqcKDB0I")
# model_name = "models/gemini-1.5-pro-latest"
model_name = "models/gemini-2.5-flash-preview-04-17"

app = Flask(__name__, static_url_path='/static')

# TEAM MEMBERS (edit this if needed)
TEAM_MEMBERS = ["Gourav Jat", "Dev shivhare", "Abhishek Rawat", "Gopal Upadhyay"]

def generate_slide_content(topic, num_slides=5, detail_level="average", style="classic", language="English"):
    detail_text = {
        "average": "moderate depth with a short paragraphs and 1 bullet points.",
        "detailed": "thorough and technical with  short paragraphs depending on context and 2 bullet points."
    }

    style_text = {
        "classic": "The presentation should have a classic, formal tone.",
        "modern": "The presentation should have a modern and visually appealing tone.",
        "dark": "The presentation should have a high-contrast, dark-themed tone."
    }

        # Prompt to Gemini
    prompt = ( f"""
You are an expert assistant skilled in generating structured PowerPoint presentations.

Generate a PowerPoint presentation on the topic: "{topic}" in the language: {language}. The total number of slides should be: {num_slides}.

Presentation requirements:
- The first slide is a title slide showing the topic and the team members: {', '.join(TEAM_MEMBERS)}.
- The last slide should be a conclusion slide with a summary of the key takeaways.
- All slides must be unique and rich in meaningful, non-repetitive content.
- Do NOT use filler lines like “This slide expands…” or “More on this topic.” Ensure every slide is informative.
- If a slide is about a process or method, include real examples or comparisons.
- Make sure there is no redundancy across slides.

Each slide (except title and conclusion) must include:
1. A slide heading.
2. A short paragraph (3–5 sentences) explaining the subtopic.
3. 3–5 meaningful bullet points that add depth, context, facts, or examples.

Format response like this (repeat for each slide):
Slide Heading: [Title of the slide]
Paragraph: [Introductory paragraph here]
- Bullet 1
- Bullet 2
- Bullet 3

Use {language} only throughout the content. The language should match the user selection exactly.

Write the content clearly, concisely, and professionally.
"""
    )

    model = genai.GenerativeModel(model_name)
    try:
        response = model.generate_content(prompt)
        if not response or not hasattr(response, "text") or not response.text.strip():
            return [{"title": "Error", "bullets": ["No content generated. Please try again."]}]

        slides_raw = response.text.strip().split("\n\n")
        slides = []

        for slide in slides_raw:
            lines = slide.strip().split("\n")
            if not lines:
                continue
            title_line = lines[0].replace("Slide Heading:", "").strip()
            bullet_points = [line.replace("-", "").strip() for line in lines[1:] if line.startswith("-")]

            if not title_line or not bullet_points:
                continue  # Skip weak slides

            slides.append({
                "title": title_line,
                "bullets": bullet_points
            })

        # Enforce number of slides and fallback
        while len(slides) < num_slides:
            slides.append({
                "title": f"Extra Slide {len(slides) + 1}",
                "bullets": [f"More insights on the topic '{topic}'."]
            })

        return slides[:num_slides]

    except Exception as e:
        return [{"title": "Error", "bullets": [str(e)]}]

def create_ppt(slide_contents, filename="AI_Presentation.pptx"):
    save_path = os.path.join(os.getcwd(), filename)
    ppt = Presentation()

    for slide_content in slide_contents:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        slide.shapes.title.text = slide_content["title"]
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()

        for point in slide_content["bullets"]:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0

    ppt.save(save_path)
    return save_path

@app.route("/")
def home():
    return render_template("index.html")

@app.route("/generate", methods=["POST"])
def generate_ppt():
    topic = request.form.get("topic")

    # Get language from form input
    language = request.form.get("language", "en")
    lang_map = {
        "en": "English",
        "fr": "French",
        "hi": "Hindi",
        "es": "Spanish",
        "de": "German"
    }
    language_name = lang_map.get(language, "English")

    num_slides_raw = request.form.get("num_slides")
    if not num_slides_raw or not num_slides_raw.isdigit():
        return jsonify({"error": "Please enter a valid number of slides."}), 400

    num_slides = int(num_slides_raw)
    detail_level = request.form.get("detail_level", "average")
    style = request.form.get("style", "classic")

    num_slides = max(1, min(250, num_slides))  # Clamp to 1–250

    # Pass language_name to the generation function
    slides_preview = generate_slide_content(topic, num_slides, detail_level, style, language_name)

    if slides_preview and slides_preview[0].get("title") == "Error":
        return jsonify({"error": slides_preview[0]["bullets"][0]}), 500

    ppt_file_path = create_ppt(slides_preview)

    return jsonify({
        "slides": slides_preview,
        "ppt_file": True
    })

@app.route("/download", methods=["GET"])
def download_ppt():
    file_path = os.path.join(os.getcwd(), "AI_Presentation.pptx")
    if not os.path.exists(file_path):
        return jsonify({"error": "PPT file not found!"}), 404
    return send_file(file_path, as_attachment=True)

if __name__ == "__main__":
    app.run(debug=True)
